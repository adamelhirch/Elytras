"""Espace de fichiers scopé (perso / projet) — mode fichier, sans base.

Cloisonnement identique à la mémoire et aux sessions : un fichier « perso » n'est
visible que par son propriétaire ; un fichier « projet » par les membres du projet.
Contenu stocké en base64 dans le filestore (limite de taille — phase 0).
"""
from __future__ import annotations

import base64
import io
import os
import time
import uuid

from . import filestore

# Décodeur HEIC/HEIF (photos iPhone) pour PIL, si le paquet est présent.
try:
    import pillow_heif  # type: ignore
    pillow_heif.register_heif_opener()
except Exception:
    pass

MAX_BYTES = int(float(os.environ.get("ELYTRAS_MAX_FILE_MB", "25")) * 1024 * 1024)   # 25 Mo par défaut
_META = ("name", "scope", "project_id", "owner_id", "size", "mime", "created_at")


def _data_dir() -> str:
    """Dossier des contenus de fichiers (sur DISQUE, hors état JSON pour ne pas l'alourdir)."""
    d = os.environ.get("ELYTRAS_FILES_DIR")
    if not d:
        base = os.path.dirname(os.path.abspath(str(filestore._PATH))) or "."
        d = os.path.join(base, "files_data")
    os.makedirs(d, exist_ok=True)
    return d


def raw_bytes(f: dict) -> bytes:
    """Octets bruts d'un fichier : depuis le disque (par id) ou le base64 hérité/en mémoire."""
    if not f:
        return b""
    if f.get("b64"):
        try:
            return base64.b64decode(f["b64"].encode())
        except Exception:
            return b""
    fid = f.get("id")
    if fid:
        p = os.path.join(_data_dir(), fid)
        if os.path.exists(p):
            try:
                with open(p, "rb") as fh:
                    return fh.read()
            except Exception:
                return b""
    return b""


def _accessible(f: dict, user_id: str, project_ids) -> bool:
    if f.get("scope") == "projet":
        return f.get("project_id") in (project_ids or [])
    return f.get("owner_id") == user_id


def list_files(user_id: str, project_ids=None) -> list[dict]:
    out = []
    for fid, f in filestore.items("files").items():
        if _accessible(f, user_id, project_ids):
            out.append({"id": fid, **{k: f.get(k) for k in _META}})
    out.sort(key=lambda x: x.get("created_at", 0), reverse=True)
    return out


def add_file(scope: str, owner_id, project_id, name: str, data_b64: str,
             mime: str = "application/octet-stream") -> str:
    raw = base64.b64decode((data_b64 or "").encode())
    if len(raw) > MAX_BYTES:
        raise ValueError(f"fichier trop volumineux (> {MAX_BYTES // (1024 * 1024)} Mo)")
    scope = "projet" if scope == "projet" else "perso"
    fid = str(uuid.uuid4())
    with open(os.path.join(_data_dir(), fid), "wb") as fh:        # contenu sur disque, hors état JSON
        fh.write(raw)
    filestore.put("files", fid, {"name": name or "fichier", "scope": scope,
                                 "project_id": project_id if scope == "projet" else None,
                                 "owner_id": None if scope == "projet" else owner_id,
                                 "size": len(raw), "mime": mime, "created_at": time.time()})
    return fid


def get_file(fid: str, user_id: str, project_ids=None) -> dict | None:
    f = filestore.items("files").get(fid)
    if not f or not _accessible(f, user_id, project_ids):
        return None
    return {"id": fid, **f}


def get_by_name(name: str, user_id: str, project_ids=None) -> dict | None:
    for fid, f in filestore.items("files").items():
        if f.get("name") == name and _accessible(f, user_id, project_ids):
            return {"id": fid, **f}
    return None


def text_of(f: dict) -> str:
    try:
        return raw_bytes(f).decode("utf-8", "replace")
    except Exception:
        return ""


# ── Extraction de contenu selon le type : PDF / Word / Excel / image (OCR) / texte ──
_TEXT_EXT = (".txt", ".csv", ".tsv", ".md", ".json", ".log", ".xml", ".html", ".htm", ".yaml", ".yml")
_IMG_EXT = (".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp", ".gif", ".heic", ".heif")


def _ocr_bytes(data: bytes) -> str:
    """OCR d'une image si Tesseract + pytesseract sont disponibles, sinon chaîne vide."""
    try:
        import pytesseract
        from PIL import Image
        im = Image.open(io.BytesIO(data))
        try:
            im = im.convert("RGB")           # normalise (HEIF, RGBA, CMYK, P…) pour Tesseract
        except Exception:
            pass
        lang = os.environ.get("OCR_LANG", "fra+eng")
        try:
            return (pytesseract.image_to_string(im, lang=lang) or "").strip()
        except Exception:
            return (pytesseract.image_to_string(im) or "").strip()      # repli langue par défaut
    except Exception:
        return ""


def _pdf_text(raw: bytes) -> str:
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(raw))
        text = "\n".join((p.extract_text() or "") for p in reader.pages).strip()
        if text:
            return text
        ocr = []                                                        # PDF scanné → OCR des images
        for p in reader.pages:
            for img in (getattr(p, "images", None) or []):
                t = _ocr_bytes(img.data)
                if t:
                    ocr.append(t)
        return "\n".join(ocr).strip() or \
            "[PDF sans couche texte — probablement scanné ; OCR indisponible ou sans résultat.]"
    except Exception as e:  # noqa: BLE001
        return f"[Échec de lecture du PDF : {e}]"


def _docx_text(raw: bytes) -> str:
    try:
        import docx
        d = docx.Document(io.BytesIO(raw))
        parts = [p.text for p in d.paragraphs]
        for table in d.tables:
            for row in table.rows:
                parts.append("\t".join(c.text for c in row.cells))
        return "\n".join(parts).strip()
    except Exception as e:  # noqa: BLE001
        return f"[Échec de lecture du Word : {e}]"


def _xlsx_text(raw: bytes) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        lines = []
        for ws in wb.worksheets:
            lines.append(f"# {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if c is None else str(c) for c in row]
                if any(cells):
                    lines.append("\t".join(cells))
        return "\n".join(lines).strip()
    except Exception as e:  # noqa: BLE001
        return f"[Échec de lecture de l'Excel : {e}]"


def _pptx_text(raw: bytes) -> str:
    try:
        import pptx
        prs = pptx.Presentation(io.BytesIO(raw))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"# Diapo {i}")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    out.append(shape.text_frame.text)
        return "\n".join(x for x in out if x).strip()
    except Exception as e:  # noqa: BLE001
        return f"[Échec de lecture du PowerPoint : {e}]"


def extract_text(f: dict) -> str:
    """Extrait le TEXTE d'un fichier selon son type. En cas d'impossibilité, renvoie un
    message clair entre crochets (jamais d'invention)."""
    if not f:
        return ""
    raw = raw_bytes(f)
    name = (f.get("name") or "").lower()
    mime = (f.get("mime") or "").lower()
    if name.endswith(_TEXT_EXT) or mime.startswith("text/") or "json" in mime:
        return raw.decode("utf-8", "replace")
    if name.endswith(".pdf") or mime == "application/pdf":
        return _pdf_text(raw)
    if name.endswith(".docx") or "wordprocessingml" in mime:
        return _docx_text(raw)
    if name.endswith((".xlsx", ".xlsm")) or "spreadsheetml" in mime:
        return _xlsx_text(raw)
    if name.endswith(".pptx") or "presentationml" in mime:
        return _pptx_text(raw)
    if name.endswith(_IMG_EXT) or mime.startswith("image/"):
        return _ocr_bytes(raw) or f"[Image « {f.get('name')} » — aucun texte détecté par l'OCR.]"
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return (f"[Fichier binaire « {f.get('name')} » ({mime or 'type inconnu'}, {len(raw)} octets) "
                "— format non pris en charge pour l'extraction de texte.]")


def delete_file(fid: str, user_id: str, project_ids=None) -> bool:
    f = filestore.items("files").get(fid)
    if not f or not _accessible(f, user_id, project_ids):
        return False
    try:
        p = os.path.join(_data_dir(), fid)
        if os.path.exists(p):
            os.remove(p)
    except Exception:
        pass
    return filestore.delete("files", fid)
