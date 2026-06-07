"""Extraction de contenu : texte, Word, Excel, image (OCR), PDF scanné (OCR), repli honnête."""
import base64
import io

import elytras.files as F


def _file(name, raw: bytes, mime=""):
    return {"name": name, "mime": mime, "b64": base64.b64encode(raw).decode()}


def _font(size=48):
    from PIL import ImageFont
    try:
        return ImageFont.truetype("DejaVuSans.ttf", size)
    except Exception:
        return ImageFont.load_default()


def test_txt():
    assert "héllo" in F.extract_text(_file("a.txt", "héllo monde".encode()))


def test_docx():
    import docx
    d = docx.Document()
    d.add_paragraph("Devis numéro 42 pour M. Dupont")
    buf = io.BytesIO()
    d.save(buf)
    assert "Devis numéro 42" in F.extract_text(_file("d.docx", buf.getvalue()))


def test_xlsx():
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"], ws["B1"], ws["A2"], ws["B2"] = "Produit", "Prix", "Vanille", 9.9
    buf = io.BytesIO()
    wb.save(buf)
    t = F.extract_text(_file("s.xlsx", buf.getvalue()))
    assert "Vanille" in t and "9.9" in t


def test_binary_fallback_no_hallucination():
    t = F.extract_text(_file("x.bin", bytes([0, 1, 2, 3, 255, 254, 7])))
    assert "non pris en charge" in t or "binaire" in t       # message honnête, pas d'invention


def test_image_ocr():
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (600, 160), "white")
    ImageDraw.Draw(img).text((20, 50), "FACTURE", fill="black", font=_font())
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    t = F.extract_text(_file("img.png", buf.getvalue(), "image/png")).upper()
    assert "FACTURE" in t


def test_scanned_pdf_ocr_or_honest():
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (600, 200), "white")
    ImageDraw.Draw(img).text((20, 60), "ETIQUETTE", fill="black", font=_font())
    buf = io.BytesIO()
    img.save(buf, format="PDF")                              # PDF image-only (= scanné)
    t = F.extract_text(_file("scan.pdf", buf.getvalue(), "application/pdf")).upper()
    # soit l'OCR récupère le texte, soit message honnête « scanné » — jamais d'invention
    assert "ETIQUETTE" in t or "PDF" in t


def test_disk_storage_roundtrip_large(state):
    blob = bytes(2 * 1024 * 1024)                          # 2 Mo (au-delà de l'ancienne limite 1 Mo)
    fid = F.add_file("perso", "u1", None, "gros.bin", base64.b64encode(blob).decode())
    meta = F.get_file(fid, "u1")
    assert meta and "b64" not in meta and meta["size"] == len(blob)   # contenu hors état JSON
    assert F.raw_bytes(meta) == blob                                  # relu depuis le disque


def test_heic_ocr():
    import pillow_heif  # noqa: F401
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (600, 160), "white")
    ImageDraw.Draw(img).text((20, 50), "DEVIS", fill="black", font=_font())
    buf = io.BytesIO()
    img.save(buf, format="HEIF")
    t = F.extract_text(_file("photo.heic", buf.getvalue(), "image/heic")).upper()
    assert "DEVIS" in t


def test_pptx():
    import pptx
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Présentation Vanille"
    buf = io.BytesIO()
    prs.save(buf)
    assert "Vanille" in F.extract_text(_file("p.pptx", buf.getvalue()))
